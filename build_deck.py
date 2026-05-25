"""One-shot pitch-deck generator. Run: python build_deck.py"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
IMG_DIR = ROOT / "images"
OUT = ROOT / "QueryDoctor_Pitch.pptx"

_shots = sorted(IMG_DIR.glob("*.png"))
if len(_shots) < 2:
    raise SystemExit(f"Expected 2 screenshots in {IMG_DIR}, found {len(_shots)}")
SHOT_DETAIL, SHOT_RESULT = _shots[0], _shots[1]
print(f"Embedding:\n  detail = {SHOT_DETAIL.name}\n  result = {SHOT_RESULT.name}")

BG = RGBColor(0x0D, 0x11, 0x17)
PANEL = RGBColor(0x16, 0x1B, 0x22)
TEXT = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)
EMERALD = RGBColor(0x10, 0xB9, 0x81)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(slide, text, left, top, width, height, *, size=18, color=TEXT,
             bold=False, align=PP_ALIGN.LEFT, font="Inter"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.name = font
        r.font.color.rgb = color
        r.font.bold = bold
    return tb


def add_accent_bar(slide, left, top, height, color=VIOLET, width=Inches(0.08)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.shadow.inherit = False
    return bar


def add_pill(slide, text, left, top, *, color=VIOLET):
    tb = slide.shapes.add_textbox(left, top, Inches(2.4), Inches(0.32))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.name = "Inter"
    r.font.bold = True
    r.font.color.rgb = color
    return tb


# ───────────────────── Slide 1 — Title ─────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)

# decorative accent strip on left
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SH)
strip.line.fill.background()
strip.fill.solid()
strip.fill.fore_color.rgb = VIOLET
strip.shadow.inherit = False

add_pill(s, "AI HACKS 2026  ·  PROBLEM STATEMENT #5", Inches(0.9), Inches(0.85))
add_text(s, "QueryDoctor", Inches(0.85), Inches(1.4), Inches(11), Inches(1.6),
         size=84, bold=True, color=TEXT)
add_text(s, "Paste a slow SQL query. Get a faster one.",
         Inches(0.9), Inches(3.4), Inches(11), Inches(0.8),
         size=26, color=MUTED)

# divider
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(5.4), Inches(2), Inches(0.03))
div.line.fill.background(); div.fill.solid(); div.fill.fore_color.rgb = VIOLET; div.shadow.inherit = False

add_text(s, "Kanishk Aman", Inches(0.9), Inches(5.6), Inches(8), Inches(0.5),
         size=18, bold=True, color=TEXT)
add_text(s, "Team LoneWolf  ·  solo build", Inches(0.9), Inches(6.0), Inches(8), Inches(0.4),
         size=14, color=MUTED)


# ───────────────────── Slide 2 — Problem ─────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)

add_pill(s, "01  ·  THE PROBLEM", Inches(0.6), Inches(0.6))
add_text(s, "Slow queries are everywhere.\nNobody reads EXPLAIN plans.",
         Inches(0.6), Inches(1.0), Inches(12), Inches(2.2),
         size=44, bold=True, color=TEXT)

# Three problem cards
cards = [
    ("$$$", "Wasted spend",
     "A single un-indexed query at production scale costs hundreds in cloud DB CPU every month."),
    ("⟳", "Latent failure",
     "Queries that pass review at 1k rows fall over at 1M. Most teams find out from a Slack alert."),
    ("¬", "Expertise gap",
     "DBAs are scarce. Junior engineers ship NOT IN anti-patterns and correlated subqueries by default."),
]
card_w = Inches(4.0); card_h = Inches(2.6); gap = Inches(0.35)
total = card_w * 3 + gap * 2
start_x = (SW - total) / 2
y = Inches(4.0)
for i, (icon, title, body) in enumerate(cards):
    x = start_x + (card_w + gap) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    box.adjustments[0] = 0.04
    box.line.color.rgb = RGBColor(0x2A, 0x33, 0x42)
    box.line.width = Pt(0.5)
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL
    box.shadow.inherit = False
    add_accent_bar(s, x + Inches(0.25), y + Inches(0.4), Inches(0.5), color=VIOLET, width=Inches(0.04))
    add_text(s, icon, x + Inches(0.4), y + Inches(0.35), Inches(1), Inches(0.6),
             size=28, bold=True, color=VIOLET)
    add_text(s, title, x + Inches(0.4), y + Inches(0.95), card_w - Inches(0.8), Inches(0.5),
             size=18, bold=True, color=TEXT)
    add_text(s, body, x + Inches(0.4), y + Inches(1.5), card_w - Inches(0.8), Inches(1.0),
             size=12, color=MUTED)


# ───────────────────── Slide 3 — Solution + demo screenshot ─────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)

add_pill(s, "02  ·  THE SOLUTION", Inches(0.6), Inches(0.5))
add_text(s, "Benchmark. Diagnose. Apply. Re-benchmark.",
         Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
         size=28, bold=True, color=TEXT)
add_text(s, "Every recommendation is verified against the real database.\n"
            "No vibes, no hand-waving — just measured milliseconds.",
         Inches(0.6), Inches(1.7), Inches(12), Inches(1.0),
         size=14, color=MUTED)

# Embed the screenshot (the 8.5× faster one)
s.shapes.add_picture(str(SHOT_RESULT), Inches(0.5), Inches(2.7), width=Inches(8.0))

# Right-side stats panel
right_x = Inches(8.8); right_w = Inches(4.0)
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, Inches(2.7),
                           right_w, Inches(4.5))
panel.adjustments[0] = 0.03
panel.line.color.rgb = VIOLET; panel.line.width = Pt(0.75)
panel.fill.solid(); panel.fill.fore_color.rgb = PANEL; panel.shadow.inherit = False

add_text(s, "DEMO RUN", right_x + Inches(0.3), Inches(2.85),
         right_w - Inches(0.6), Inches(0.3), size=10, bold=True, color=MUTED)

add_text(s, "Best-selling products", right_x + Inches(0.3), Inches(3.35),
         right_w - Inches(0.6), Inches(0.4), size=13, color=TEXT)
add_text(s, "aggregate over 260k rows", right_x + Inches(0.3), Inches(3.7),
         right_w - Inches(0.6), Inches(0.3), size=10, color=MUTED)

# Baseline
add_text(s, "Baseline", right_x + Inches(0.3), Inches(4.15),
         right_w - Inches(0.6), Inches(0.3), size=10, color=MUTED)
add_text(s, "179.92 ms", right_x + Inches(0.3), Inches(4.4),
         right_w - Inches(0.6), Inches(0.5), size=24, bold=True, color=ROSE)

# After
add_text(s, "After fix", right_x + Inches(0.3), Inches(5.05),
         right_w - Inches(0.6), Inches(0.3), size=10, color=MUTED)
add_text(s, "21.26 ms", right_x + Inches(0.3), Inches(5.3),
         right_w - Inches(0.6), Inches(0.5), size=24, bold=True, color=EMERALD)

# Speedup
div2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x + Inches(0.3), Inches(6.0),
                          right_w - Inches(0.6), Inches(0.02))
div2.line.fill.background(); div2.fill.solid(); div2.fill.fore_color.rgb = VIOLET
div2.shadow.inherit = False

add_text(s, "8.5× faster", right_x + Inches(0.3), Inches(6.15),
         right_w - Inches(0.6), Inches(0.6), size=22, bold=True, color=VIOLET)
add_text(s, "with one click", right_x + Inches(0.3), Inches(6.65),
         right_w - Inches(0.6), Inches(0.3), size=11, color=MUTED)


# ───────────────────── Slide 4 — How it works ─────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)

add_pill(s, "03  ·  HOW IT WORKS", Inches(0.6), Inches(0.5))
add_text(s, "Grounded suggestions, not generic SQL tips.",
         Inches(0.6), Inches(0.9), Inches(12), Inches(0.7),
         size=26, bold=True, color=TEXT)

# Pipeline arrows row
steps = [
    ("Query +\nschema", VIOLET),
    ("EXPLAIN\nplan", VIOLET),
    ("Gemini\n2.5 Flash", VIOLET),
    ("Diagnosis +\nfix SQL", VIOLET),
    ("Apply &\nre-benchmark", EMERALD),
]
step_w = Inches(2.1); step_h = Inches(1.05); step_gap = Inches(0.12)
total_w = step_w * len(steps) + step_gap * (len(steps) - 1)
y = Inches(1.75)
sx = (SW - total_w) / 2
for i, (label, col) in enumerate(steps):
    x = sx + (step_w + step_gap) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, step_w, step_h)
    box.adjustments[0] = 0.15
    box.line.color.rgb = col; box.line.width = Pt(0.75)
    box.fill.solid(); box.fill.fore_color.rgb = PANEL; box.shadow.inherit = False
    add_text(s, label, x, y + Inches(0.18), step_w, step_h - Inches(0.3),
             size=12, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

# Embed detail screenshot
s.shapes.add_picture(str(SHOT_DETAIL), Inches(0.5), Inches(3.05), width=Inches(8.0))

# Right-side bullets
rx = Inches(8.8); rw = Inches(4.1)
add_text(s, "Why this is different", rx, Inches(3.05), rw, Inches(0.4),
         size=14, bold=True, color=VIOLET)

bullets = [
    ("Schema-aware",
     "Gemini sees actual tables, existing indexes, and row counts — recommendations fit your database, not a textbook."),
    ("Plan-grounded",
     "Suggestions are checked against the EXPLAIN plan: a fix is only proposed if it converts SCAN → SEARCH."),
    ("Anti-pattern aware",
     "Catches NOT IN, correlated subqueries, and missing covering indexes. Rewrites SQL, not just adds indexes."),
    ("Verified, not guessed",
     "Every fix is applied and re-run on the live database. The speedup number is measured, not promised."),
]
by = Inches(3.6)
for title, body in bullets:
    add_accent_bar(s, rx, by + Inches(0.05), Inches(0.8), color=VIOLET, width=Inches(0.04))
    add_text(s, title, rx + Inches(0.15), by, rw - Inches(0.2), Inches(0.35),
             size=12, bold=True, color=TEXT)
    add_text(s, body, rx + Inches(0.15), by + Inches(0.32), rw - Inches(0.2), Inches(0.7),
             size=9, color=MUTED)
    by += Inches(0.9)


# ───────────────────── Slide 5 — Impact + Tech + Close ─────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)

add_pill(s, "04  ·  IMPACT  &  STACK", Inches(0.6), Inches(0.5))
add_text(s, "One click. Real dollars.", Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
         size=30, bold=True, color=TEXT)

# Big impact metrics
metrics = [
    ("158.66 ms", "saved per request", ROSE),
    ("380 min", "DB-CPU saved / day", VIOLET),
    ("$38.08", "monthly cost saving", EMERALD),
]
mw = Inches(3.6); mh = Inches(2.0); mg = Inches(0.3)
total_mw = mw * 3 + mg * 2
mx = (SW - total_mw) / 2
my = Inches(2.0)
for i, (big, label, col) in enumerate(metrics):
    x = mx + (mw + mg) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, my, mw, mh)
    box.adjustments[0] = 0.08
    box.line.color.rgb = col; box.line.width = Pt(0.75)
    box.fill.solid(); box.fill.fore_color.rgb = PANEL; box.shadow.inherit = False
    add_text(s, big, x, my + Inches(0.35), mw, Inches(0.9),
             size=44, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, label, x, my + Inches(1.35), mw, Inches(0.4),
             size=13, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s, "Modeled at 100 req/min on a single query. Multiply across a real codebase.",
         Inches(0.6), Inches(4.15), Inches(12), Inches(0.4),
         size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Tech stack row
add_text(s, "BUILT WITH", Inches(0.6), Inches(4.85), Inches(6), Inches(0.3),
         size=11, bold=True, color=MUTED)
add_text(s, "Python  ·  Streamlit  ·  SQLite  ·  Plotly  ·  Google Gemini 2.5 Flash",
         Inches(0.6), Inches(5.15), Inches(9), Inches(0.5),
         size=16, color=TEXT)

add_text(s, "WHAT'S NEXT", Inches(0.6), Inches(5.9), Inches(6), Inches(0.3),
         size=11, bold=True, color=MUTED)
add_text(s, "Postgres + EXPLAIN ANALYZE  ·  Pre-commit SQL lint  ·  Fix history across the codebase",
         Inches(0.6), Inches(6.2), Inches(12), Inches(0.5),
         size=13, color=TEXT)

# Footer
foot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), SW, Inches(0.5))
foot.line.fill.background(); foot.fill.solid(); foot.fill.fore_color.rgb = PANEL
foot.shadow.inherit = False
add_text(s, "Kanishk Aman  ·  Team LoneWolf  ·  AI Hacks 2026",
         Inches(0.6), Inches(7.12), Inches(12), Inches(0.3),
         size=11, color=MUTED)
add_text(s, "QueryDoctor", SW - Inches(2.6), Inches(7.12), Inches(2), Inches(0.3),
         size=11, bold=True, color=VIOLET, align=PP_ALIGN.RIGHT)


prs.save(str(OUT))
print(f"Wrote {OUT}  ({OUT.stat().st_size / 1024:.1f} KB)")
